from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import copy

# === Цветовая палитра ===
COLOR_BG_DARK = RGBColor(0x0D, 0x18, 0x2D)
COLOR_BG_LIGHT = RGBColor(0xF0, 0xF3, 0xF9)
COLOR_ACCENT_BLUE = RGBColor(0x27, 0x7F, 0xD6)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_GRAY_TEXT = RGBColor(0x55, 0x5B, 0x62)
COLOR_TITLE_TEXT = RGBColor(0x1C, 0x2A, 0x45)
COLOR_RED = RGBColor(0xE5, 0x3D, 0x3D)
COLOR_ORANGE = RGBColor(0xF4, 0xA2, 0x2A)
COLOR_GREEN = RGBColor(0x3A, 0x85, 0x53)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# === Вспомогательные функции ===
def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_to_shape(shape, text, font_size, bold=False, align=PP_ALIGN.LEFT, color=COLOR_GRAY_TEXT, margin=0.1):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return p

# === СЛАЙД 1 (Титульный) ===
def create_slide_1():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_DARK)
    
    title_shape = add_shape(slide, Inches(1.0), Inches(1.5), Inches(12), Inches(2.0), COLOR_BG_DARK)
    add_text_to_shape(title_shape, "Интеграция сервисов Mail в VK", 52, bold=True, color=COLOR_WHITE, margin=0.4)
    sub_shape = add_shape(slide, Inches(1.0), Inches(3.2), Inches(12), Inches(1.0), COLOR_BG_DARK)
    add_text_to_shape(sub_shape, "Потенциал аудитории, финансовая модель и MVP внедрения", 24, bold=False, color=RGBColor(0xCC, 0xDD, 0xEE), margin=0.4)

    circle_outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(0.5), Inches(2.5), Inches(2.5))
    circle_outer.fill.solid()
    circle_outer.fill.fore_color.rgb = RGBColor(0x1D, 0x37, 0x5A)
    circle_outer.line.fill.background()
    circle_inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.0), Inches(1.0), Inches(1.5), Inches(1.5))
    circle_inner.fill.solid()
    circle_inner.fill.fore_color.rgb = COLOR_ACCENT_BLUE
    circle_inner.line.fill.background()
    vk_text = slide.shapes.add_textbox(Inches(11.3), Inches(1.1), Inches(1.5), Inches(1.5))
    p = vk_text.text_frame.paragraphs[0]
    p.text = "VK"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

# === СЛАЙД 2 (Контекст) ===
def create_slide_2():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_LIGHT)
    header = add_shape(slide, Inches(0.8), Inches(2.0), Inches(12), Inches(3.85), COLOR_BG_LIGHT) 
    add_text_to_shape(header, "01 Контекст: Борьба за 'Липкость'", 40, bold=True, color=COLOR_TITLE_TEXT)
    ns_block = add_shape(slide, Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.8), COLOR_ACCENT_BLUE)
    add_text_to_shape(ns_block, "NORTH STAR  •  DAU / MAU — Stickiness (Доля ежедневных пользователей)", 16, bold=True, color=COLOR_WHITE)
    cards = [
        {'title': '🎯 Цель', 'text': 'Выявить ключевые барьеры, мешающие возвращаться в VK и взаимодействовать с контентом.'},
        {'title': '📉 Важность', 'text': 'Снижение DAU/MAU = падение рекламной выручки (меньше показов).'},
        {'title': '⚔️ Конкуренты', 'text': 'Удержание падает на фоне Telegram и YouTube. Их DAU/MAU растут.'}
    ]
    for i, card_data in enumerate(cards):
        x = Inches(0.8) + (i * (Inches(3.7) + Inches(0.2)))
        y = Inches(2.8)
        card = add_shape(slide, x, y, Inches(3.7), Inches(3.0), COLOR_WHITE, COLOR_WHITE)
        title_box = add_shape(slide, x, y, Inches(3.7), Inches(0.6), COLOR_WHITE)
        add_text_to_shape(title_box, card_data['title'], 16, bold=True, color=COLOR_TITLE_TEXT)
        desc_box = add_shape(slide, x, Inches(3.4), Inches(3.7), Inches(2.4), COLOR_WHITE)
        add_text_to_shape(desc_box, card_data['text'], 16, color=COLOR_GRAY_TEXT)

# === СЛАЙД 3 (Приоритет продуктов) ===
def create_slide_3():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_LIGHT)
    add_text_to_shape(add_shape(slide, Inches(0.8), Inches(0.5), Inches(12), Inches(1.0), COLOR_BG_LIGHT), "02 Приоритетные сервисы Mail", 38, bold=True, color=COLOR_TITLE_TEXT)
    products = [
        ("1 место", "Задачи", "Виджет в чате. Создание из сообщения. Привычный паттерн."),
        ("1 место", "Календарь", "Встречи прямо из переписки. Снижение когнитивной нагрузки."),
        ("1 место", "Заметки", "Аналог 'Избранного'. 100+ млн пользователей."),
        ("1 место", "VK Звонки", "Расширение текущих звонков. Без обучения."),
        ("2 место", "Облако", "Выгодные тарифы (≈10% дешевле) Яндекс/Google."),
        ("3 место", "Почта", "Сложный функционал, долгая интеграция.")
    ]
    for i, (rank, name, desc) in enumerate(products):
        x = Inches(0.8) if i < 3 else Inches(6.8)
        y = Inches(1.8) + ((i % 3) * Inches(1.5))
        card = add_shape(slide, x, y, Inches(5.5), Inches(1.4), COLOR_WHITE, COLOR_WHITE)
        add_text_to_shape(card, f"{rank} • {name}\n{desc}", 13, bold=False, color=COLOR_GRAY_TEXT)

# === СЛАЙД 4 (Финансовый потенциал) ===
def create_slide_4():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_LIGHT)
    add_text_to_shape(add_shape(slide, Inches(0.8), Inches(0.5), Inches(12), Inches(1.0), COLOR_BG_LIGHT), "03 Финансовый потенциал (Годовая выручка)", 38, bold=True, color=COLOR_TITLE_TEXT)
    
    data = [
        ("Задачи", "20%", "21.4 млн", "570 ₽", "12.2 млрд ₽"),
        ("Календарь", "25%", "26.75 млн", "570 ₽", "15.2 млрд ₽"),
        ("Заметки", "20%", "21.4 млн", "570 ₽", "12.2 млрд ₽"),
        ("VK Звонки", "50%", "53.5 млн", "570 ₽", "30.5 млрд ₽"),
        ("Облако", "30%", "32.1 млн", "570 ₽", "18.3 млрд ₽"),
    ]
    
    y_offset = Inches(1.8)
    for i, (name, pen, aud, arpu, rev) in enumerate(data):
        y = y_offset + i * Inches(0.85)
        card = add_shape(slide, Inches(1.5), y, Inches(10.0), Inches(0.75), COLOR_WHITE, COLOR_WHITE)
        add_text_to_shape(card, f"{name} | PR: {pen} | Аудитория: {aud} | ARPU: {arpu}/год | Выручка: {rev}", 13, bold=True, color=COLOR_GRAY_TEXT)
    
    total = add_shape(slide, Inches(1.5), y_offset + 5 * Inches(0.85), Inches(10.0), Inches(0.8), COLOR_ACCENT_BLUE)
    add_text_to_shape(total, "ИТОГО: 88.4 млрд руб./год (∑ всех сервисов)", 16, bold=True, color=COLOR_WHITE)

# === СЛАЙД 5 (Конкуренты) ===
def create_slide_5():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_LIGHT)
    add_text_to_shape(add_shape(slide, Inches(0.8), Inches(0.5), Inches(12), Inches(1.0), COLOR_BG_LIGHT), "04 Конкурентный анализ (Россия)", 38, bold=True, color=COLOR_TITLE_TEXT)
    
    comps = [
        ("Базовые чаты", "✓", "✓", "✓", "✓"),
        ("Видеозвонки", "✓", "✓", "✓", "✓"),
        ("Сторис", "✓", "✓", "✓", "✓"),
        ("Приватные каналы", "✓", "✓", "✗", "✗"),
        ("Кошелек/Переводы", "✓", "✓", "✗", "✗"),
        ("Избранное", "✓", "✓", "✗", "✓"),
        ("Вложенные каналы", "✓", "✓", "✗", "✗")
    ]
    
    y = Inches(1.8)
    for i, (func, m, t, tm, tc) in enumerate(comps):
        card = add_shape(slide, Inches(0.8), y + i * Inches(0.7), Inches(11.5), Inches(0.6), COLOR_WHITE, COLOR_WHITE)
        add_text_to_shape(card, f"{func: <35}  |  MAX: {m} |  Telegram: {t} |  TamTam: {tm} |  TenChat: {tc}", 13, bold=False, color=COLOR_GRAY_TEXT)

# === СЛАЙД 6 (Гипотезы и метрики) ===
def create_slide_6():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_LIGHT)
    add_text_to_shape(add_shape(slide, Inches(0.8), Inches(0.5), Inches(12), Inches(1.0), COLOR_BG_LIGHT), "05 Гипотезы и метрики", 38, bold=True, color=COLOR_TITLE_TEXT)
    
    h = [
        ("🔴 Высокий приоритет", "• Stories → +10% retention, -14% bounce\n• Календарь → +10% scroll depth\n• Заметки → +5% retention\n• Задачи → +20% DAU\n• VK Звонки → +5% retention\nПроверка: A/B тест (2 нед.) на 50%"),
        ("🟠 Средний приоритет", "• Создание ботов → +10-15% conversion\n• API Integration → +10% retention\nПроверка: Колич. опрос (500+) + анализ запросов."),
        ("🟢 Низкий приоритет", "• Приватные каналы → +8% retention, +12% DAU\n• Кошелек → +5% retention, +9% DAU\nПроверка: Глубинные интервью (10-15 продавцов)")
    ]
    for i, (priority, text) in enumerate(h):
        x = Inches(0.8) + i * (Inches(3.8) + Inches(0.2))
        y = Inches(1.8)
        grp = add_shape(slide, x, y, Inches(3.8), Inches(5.2), COLOR_WHITE, COLOR_WHITE)
        add_text_to_shape(grp, f"{priority}\n\n{text}", 13, bold=False, color=COLOR_GRAY_TEXT)

# === СЛАЙД 7 (MVP и Discovery) ===
def create_slide_7():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_LIGHT)
    add_text_to_shape(add_shape(slide, Inches(0.8), Inches(0.5), Inches(12), Inches(1.0), COLOR_BG_LIGHT), "06 MVP + Discovery", 38, bold=True, color=COLOR_TITLE_TEXT)
    
    col1 = add_shape(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5), COLOR_WHITE, COLOR_WHITE)
    add_text_to_shape(col1, "🔬 Discovery до этапа разработки:\n\n• Custdev (10-15 польз.)\n• Опрос внутри приложения (1000+)\n• Анализ точек выхода (Яндекс.Метрика)", 16, bold=True, color=COLOR_GRAY_TEXT)
    
    col2 = add_shape(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5), COLOR_WHITE, COLOR_WHITE)
    add_text_to_shape(col2, "🚀 MVP: Заметки + Задачи как единый виджет в чате\n\nПочему: \n• Мин. сложность разработки\n• Закрывает боль 'закрепленных сообщений'\n• Быстрый эффект на retention\n\nФункционал: \n• Создание из сообщения\n• Дедлайн + напоминание\n• Вкладка 'Все задачи'\n• Шеринг в чате", 16, bold=True, color=COLOR_GRAY_TEXT)

# === СЛАЙД 8 (Дорожная карта) ===
def create_slide_8():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_LIGHT)
    add_text_to_shape(add_shape(slide, Inches(0.8), Inches(0.5), Inches(12), Inches(1.0), COLOR_BG_LIGHT), "07 Дорожная карта внедрения", 38, bold=True, color=COLOR_TITLE_TEXT)
    
    waves = [
        ("Wave 1 (0-3 мес.)", "MVP", "Заметки + Задачи (виджет в чате)", COLOR_ACCENT_BLUE),
        ("Wave 2 (3-6 мес.)", "Расширение", "Календарь + VK Звонки", COLOR_ORANGE),
        ("Wave 3 (6-12 мес.)", "Масштабирование", "Облако + Почта (фоново)", RGBColor(0x3A, 0x85, 0x53))
    ]
    
    y = Inches(2.0)
    for i, (wave, stage, description, color) in enumerate(waves):
        # Полоса этапа
        bar = add_shape(slide, Inches(2.0), y + i * Inches(1.5), Inches(9.0), Inches(0.8), color)
        add_text_to_shape(bar, f"{wave}  •  {stage}\n{description}", 16, bold=True, color=COLOR_WHITE)
        
        # Стрелка между этапами
        if i < 2:
            arrow = add_shape(slide, Inches(6.0), y + i * Inches(1.5) + Inches(0.8), Inches(1.0), Inches(0.2), COLOR_ACCENT_BLUE) # Упрощенная стрелка

# === СЛАЙД 9 (Риски и зависимости) ===
def create_slide_9():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_LIGHT)
    add_text_to_shape(add_shape(slide, Inches(0.8), Inches(0.5), Inches(12), Inches(1.0), COLOR_BG_LIGHT), "08 Риски и зависимости", 38, bold=True, color=COLOR_TITLE_TEXT)
    
    risks = [
        ("🔴 Технические", "• Сложность интеграции Почты\n• Зависимость от API VK Мессенджер\n• Качество связи при звонках"),
        ("🟠 Рыночные", "• Высокая активность Telegram\n• Уход пользователей при технических сбоях\n• Низкий Penetration Rate у Облака"),
        ("🟢 Организационные", "• Недостаток ресурсов разработки\n• Рассинхрон с командой VK MAX\n• Приоритеты других продуктов VK")
    ]
    
    for i, (category, text) in enumerate(risks):
        x = Inches(0.8) + i * (Inches(3.8) + Inches(0.2))
        y = Inches(2.0)
        card = add_shape(slide, x, y, Inches(3.8), Inches(4.0), COLOR_WHITE, COLOR_WHITE)
        add_text_to_shape(card, f"{category}\n\n{text}", 14, bold=False, color=COLOR_GRAY_TEXT)

# === СЛАЙД 10 (Итоги и следующий шаг) ===
def create_slide_10():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_DARK)
    
    summary = add_shape(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(3.5), COLOR_BG_DARK)
    add_text_to_shape(summary, "09 Итоги и следующий шаг", 40, bold=True, color=COLOR_WHITE)
    
    lines = [
        "💰 Потенциальная выручка от внедрения 5 сервисов > 88 млрд руб./год",
        "🚀 MVP: Заметки + Задачи (закрывает боль 'закрепленных сообщений')",
        "📈 Дорожная карта на 12 месяцев: от MVP к масштабированию",
        "⚠️ Ключевой риск: конкуренция с Telegram и техническая сложность интеграции",
        "✅ Решение: запустить A/B тест и Custdev в ближайшие 2 недели"
    ]
    
    y = Inches(3.0)
    for line in lines:
        card = add_shape(slide, Inches(1.5), y, Inches(10.0), Inches(0.6), COLOR_WHITE, COLOR_WHITE)
        add_text_to_shape(card, line, 18, bold=True, color=COLOR_GRAY_TEXT)
        y += Inches(0.8)
    
    action = add_shape(slide, Inches(1.5), y + Inches(0.5), Inches(10.0), Inches(0.8), COLOR_ACCENT_BLUE)
    add_text_to_shape(action, "👉 Рекомендация: Немедленный старт MVP (Заметки+Задачи), запуск A/B теста на 50%", 16, bold=True, color=COLOR_WHITE)

# === Сборка всех слайдов ===
create_slide_1()
create_slide_2()
create_slide_3()
create_slide_4()
create_slide_5()
create_slide_6()
create_slide_7()
create_slide_8()
create_slide_9()
create_slide_10()

prs.save("VK_Mail_Case_10_Slides.pptx")
print("✅ Презентация создана: VK_Mail_Case_10_Slides.pptx")
